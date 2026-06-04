"""Build the SUML presentation deck (slides.pptx). Run: python docs/build_slides.py"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x02, 0x80, 0x90)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
LIGHT = RGBColor(0xF5, 0xF7, 0xFB)
DARK = RGBColor(0x22, 0x22, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5A, 0x63, 0x78)
HEAD = "Trebuchet MS"
BODY = "Calibri"
IMG = Path("docs/img")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=LIGHT, motif=True):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    if motif:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()
        bar.shadow.inherit = False
    return s


def text(
    s,
    left,
    top,
    width,
    height,
    lines,
    size=16,
    color=DARK,
    bold=False,
    font=BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    bullet=False,
    space=8,
):
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(space)
        run = para.add_run()
        run.text = ("•  " + line) if bullet else line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return box


def title(s, label, color=NAVY):
    text(s, 0.7, 0.5, 12.0, 1.0, label, size=32, color=color, bold=True, font=HEAD)


def rrect(s, left, top, width, height, fill, line=None):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def card(s, left, top, width, height, big, label):
    box = rrect(s, left, top, width, height, WHITE, line=ICE)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = big
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.name = HEAD
    r.font.color.rgb = TEAL
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(13)
    r2.font.name = BODY
    r2.font.color.rgb = MUTED


def panel(s, left, top, width, height, header, lines):
    box = rrect(s, left, top, width, height, WHITE, line=ICE)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(12)
    tf.margin_bottom = Pt(12)
    p = tf.paragraphs[0]
    rh = p.add_run()
    rh.text = header
    rh.font.size = Pt(20)
    rh.font.bold = True
    rh.font.name = HEAD
    rh.font.color.rgb = NAVY
    p.space_after = Pt(8)
    for line in lines:
        para = tf.add_paragraph()
        para.space_after = Pt(6)
        run = para.add_run()
        run.text = "•  " + line
        run.font.size = Pt(15)
        run.font.name = BODY
        run.font.color.rgb = DARK


# --- Slide 1: title (dark) ---
s = slide(NAVY, motif=False)
rrect(s, 0, 0, 0.32, 7.5, TEAL)
text(s, 0.9, 2.3, 11.5, 1.4, "Food Delivery ETA", size=54, color=WHITE, bold=True, font=HEAD)
text(
    s,
    0.9,
    3.7,
    11.5,
    0.9,
    "Predykcja czasu dostawy jedzenia — AutoML, FastAPI, Streamlit, Docker",
    size=22,
    color=ICE,
    font=BODY,
)
text(
    s,
    0.9,
    5.5,
    11.5,
    0.5,
    "SUML — Środowiska uruchomieniowe ML · PJATK",
    size=16,
    color=WHITE,
    bold=True,
    font=BODY,
)
text(s, 0.9, 6.0, 11.5, 0.5, "Grupa: [numery indeksów]", size=15, color=ICE, font=BODY)

# --- Slide 2: problem ---
s = slide()
title(s, "Problem i wartość biznesowa")
text(
    s,
    0.7,
    1.9,
    6.6,
    4.6,
    [
        "Platformy dostaw (Wolt, Glovo, Uber Eats) potrzebują trafnego ETA już w momencie zamówienia.",
        "Dobre ETA = realne oczekiwania klienta, sprawniejsza dyspozycja kurierów i wczesne flagowanie opóźnień.",
        "Aplikacja przewiduje czas dostawy (regresja) na podstawie cech zamówienia, trasy i kontekstu.",
    ],
    size=17,
    bullet=True,
    space=14,
)
s_img = IMG / "distance_vs_time.png"
if s_img.exists():
    s.shapes.add_picture(str(s_img), Inches(7.6), Inches(2.1), width=Inches(5.2))

# --- Slide 3: data ---
s = slide()
title(s, "Dane")
text(
    s,
    0.7,
    1.9,
    6.6,
    4.6,
    [
        "Źródło: Kaggle denkuznetz/food-delivery-time-prediction (1000 wierszy, 9 kolumn).",
        "CSV dołączony do repo + generator danych syntetycznych o tym samym schemacie (fallback).",
        "Braki: po 30 w 4 kolumnach → imputacja w pipeline (mediana / najczęstsza wartość).",
        "Cel: czas dostawy 8–153 min (średnio ~57).",
    ],
    size=16,
    bullet=True,
    space=12,
)
t_img = IMG / "target_hist.png"
if t_img.exists():
    s.shapes.add_picture(str(t_img), Inches(7.7), Inches(2.2), width=Inches(5.0))

# --- Slide 4: architecture ---
s = slide()
title(s, "Architektura: data | model | app")
boxes = [
    ("data", "load + walidacja\nsplit + preprocessing"),
    ("model", "FLAML AutoML\nmodel.joblib + metrics.json"),
    ("app", "FastAPI /predict\nStreamlit UI"),
]
xs = [0.9, 4.9, 8.9]
for (name, desc), x in zip(boxes, xs):
    b = rrect(s, x, 2.6, 3.5, 2.0, ICE)
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rr = p.add_run()
    rr.text = name
    rr.font.size = Pt(26)
    rr.font.bold = True
    rr.font.name = HEAD
    rr.font.color.rgb = NAVY
    for ln in desc.split("\n"):
        pp = tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        r3 = pp.add_run()
        r3.text = ln
        r3.font.size = Pt(13)
        r3.font.name = BODY
        r3.font.color.rgb = DARK
for ax in (4.45, 8.45):
    arr = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(3.35), Inches(0.4), Inches(0.5)
    )
    arr.fill.solid()
    arr.fill.fore_color.rgb = TEAL
    arr.line.fill.background()
    arr.shadow.inherit = False
banner = rrect(s, 0.9, 5.1, 11.5, 1.0, NAVY)
tf = banner.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Pt(16)
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Wszystko sterowane przez config.yaml — zmiana datasetu lub retuning to zmiana configu, nie kodu."
r.font.size = Pt(16)
r.font.bold = True
r.font.name = BODY
r.font.color.rgb = WHITE

# --- Slide 5: AutoML ---
s = slide()
title(s, "Model: AutoML (FLAML)")
text(
    s,
    0.7,
    1.9,
    6.6,
    4.6,
    [
        "AutoML.fit przeszukuje estymatory (lgbm, rf, extra_tree) i składa je w ensemble (stacking).",
        "Cały preprocessing + model zapisany jako jeden sklearn Pipeline (model.joblib).",
        "Imputacja + one-hot w pipeline → spójność trening / serwowanie, bez wycieków.",
        "Konfiguracja AutoML w config.yaml (budżet, metryka, lista estymatorów).",
    ],
    size=16,
    bullet=True,
    space=12,
)
f_img = IMG / "feat_importance.png"
if f_img.exists():
    s.shapes.add_picture(str(f_img), Inches(7.6), Inches(2.2), width=Inches(5.2))

# --- Slide 6: results ---
s = slide()
title(s, "Wyniki modelu")
card(s, 0.9, 2.5, 3.6, 2.2, "6,1", "MAE (minuty)")
card(s, 4.85, 2.5, 3.6, 2.2, "9,0", "RMSE")
card(s, 8.8, 2.5, 3.6, 2.2, "0,82", "R²")
text(
    s,
    0.9,
    5.1,
    11.5,
    0.8,
    "Model: FLAML ensemble (baza LightGBM) · ocena na 20% holdout · najważniejsza cecha: dystans.",
    size=16,
    color=MUTED,
    align=PP_ALIGN.CENTER,
)

# --- Slide 7: serving ---
s = slide()
title(s, "Wystawienie: API + UI")
panel(
    s,
    0.8,
    2.0,
    5.7,
    4.2,
    "FastAPI (usługa)",
    [
        "POST /predict — predykcja ETA (minuty)",
        "GET /health — status + czy model wczytany",
        "GET /model-info — metryki i metadane",
        "Interaktywne /docs (OpenAPI) za darmo",
        "Walidacja wejścia Pydantic → 422 dla błędnych danych",
    ],
)
panel(
    s,
    6.8,
    2.0,
    5.7,
    4.2,
    "Streamlit (UI)",
    [
        "Formularz cech zamówienia",
        "Woła API i pokazuje przewidziany czas",
        "Panel z metrykami modelu",
        "Wykres ważności cech",
        "Konfigurowalny adres API (env / config)",
    ],
)

# --- Slide 8: portability ---
s = slide()
title(s, "Przenoszalność i odtwarzalność")
text(
    s,
    0.7,
    1.9,
    7.2,
    4.6,
    [
        "docker compose up --build → dwa serwisy: api (FastAPI) + ui (Streamlit).",
        "Model trenowany podczas budowania obrazu — gotowy przy pierwszym żądaniu.",
        "Zależności przypięte (requirements.txt); obraz python:3.11-slim, użytkownik non-root.",
        "Zero konfiguracji systemu operacyjnego; .gitattributes wymusza końce linii LF.",
        "config.yaml jako jedno źródło prawdy.",
    ],
    size=16,
    bullet=True,
    space=11,
)
banner = rrect(s, 8.3, 2.4, 4.2, 2.4, TEAL)
tf = banner.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Pt(14)
tf.margin_right = Pt(14)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Jedna komenda\nuruchamia całość"
r.font.size = Pt(24)
r.font.bold = True
r.font.name = HEAD
r.font.color.rgb = WHITE

# --- Slide 9: quality ---
s = slide()
title(s, "Jakość kodu i organizacja")
card(s, 0.9, 2.4, 3.6, 2.0, "10/10", "pylint")
card(s, 4.85, 2.4, 3.6, 2.0, "26", "testów (pytest)")
card(s, 8.8, 2.4, 3.6, 2.0, "CI", "GitHub Actions")
text(
    s,
    0.7,
    4.8,
    12.0,
    1.8,
    [
        "PEP8 + black + isort, type hints i docstringi w całym kodzie.",
        "Ścisły podział data | model | app; modularność na poziomie funkcji i katalogów.",
        "Czytelna historia commitów; spec i plan w docs/.",
    ],
    size=16,
    bullet=True,
    space=10,
)

# --- Slide 10: demo + closing (dark) ---
s = slide(NAVY, motif=False)
rrect(s, 0, 0, 0.32, 7.5, TEAL)
title(s, "Demo i podsumowanie", color=WHITE)
text(
    s,
    0.9,
    1.9,
    11.4,
    4.0,
    [
        "Demo: docker compose up --build → UI :8501, API :8000/docs → predykcja na żywo.",
        "Retraining: nowy CSV w data/raw/ + python -m model.train (bez zmian w kodzie).",
        "Repo: github.com/nimzoi/SUML-project",
    ],
    size=18,
    color=ICE,
    bullet=True,
    space=14,
)
text(s, 0.9, 6.0, 11.4, 0.8, "Dziękujemy — pytania?", size=24, color=WHITE, bold=True, font=HEAD)

prs.save("slides.pptx")
print("saved slides.pptx with", len(prs.slides), "slides")
